import io
import os
from typing import IO, Any, Dict

from pptx import Presentation

from wisup_e2m.utils.image_util import has_transparent_background


def get_pptx_images(
    file_name: str,
    file: IO[bytes],
    target_image_dir: str,
    ignore_transparent_images: bool = False,
) -> Dict[int, Any]:
    os.makedirs(target_image_dir, exist_ok=True)
    image_dict = {}
    prs = Presentation(file_name or file)
    for idx, slide in enumerate(prs.slides):
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:
                image = shape.image
                image_bytes: IO[bytes] = image.blob

                if ignore_transparent_images and has_transparent_background(
                    io.BytesIO(image_bytes)
                ):
                    continue

                image_name = f"{idx}_{image_count}.{image.ext}"
                image_file = os.path.join(target_image_dir, image_name)
                with open(image_file, "wb") as f:
                    f.write(image_bytes)
                if idx not in image_dict:
                    image_dict[idx] = []
                image_dict[idx].append(
                    {
                        "slide_number": idx,
                        "image_file": image_file,
                        "image_name": image_name,
                    }
                )
                image_count += 1

    # {page_number: [{'slide_number': 0, 'image_file': 'extracted_images/0_0.png', 'image_name': '0_0.png'}]}
    return image_dict
